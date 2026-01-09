from pptx import Presentation
import io

def extract_pptx(file_bytes):
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        text_content = []
        
        for i, slide in enumerate(prs.slides):
            text_content.append(f"### Slide {i+1}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        return "\n\n".join(text_content)
    except Exception as e:
        return f"[Error extracting PPTX: {str(e)}]"
